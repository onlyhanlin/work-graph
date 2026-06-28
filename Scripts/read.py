import os
import subprocess
import sys
from utils import (
    create_markdown_with_metadata, 
    create_temp_file, 
    delete_temp_file, 
    run_markitdown,
    get_doc_dir,
    ensure_dir_exists
)

IMAGE_EXTENSIONS = ['.png', '.jpg', '.jpeg', '.bmp']
TEXT_EXTENSIONS = ['.txt', '.md']


def extract_text_from_image(image_path):
    try:
        from rapidocr_onnxruntime import RapidOCR
        
        ocr = RapidOCR()
        result, _ = ocr(image_path)
        
        if result is None:
            return ""
        
        text_lines = []
        for line in result:
            if line[1] and isinstance(line[1], str):
                text_lines.append(line[1])
        
        return '\n'.join(text_lines)
    except ImportError:
        print("✗ 未安装 rapidocr-onnxruntime")
        return ""
    except Exception as e:
        print(f"✗ OCR识别失败: {e}")
        return ""


def read_pdf(file_path):
    try:
        import fitz
        
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    except ImportError:
        print("✗ 未安装 PyMuPDF")
        return None
    except Exception as e:
        print(f"✗ 读取PDF失败: {e}")
        return None


def read_docx(file_path):
    try:
        from docx import Document
        
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + '\n'
        return text
    except ImportError:
        print("✗ 未安装 python-docx")
        return None
    except Exception as e:
        print(f"✗ 读取DOCX失败: {e}")
        return None


def read_pptx(file_path):
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        return text
    except ImportError:
        print("✗ 未安装 python-pptx")
        return None
    except Exception as e:
        print(f"✗ 读取PPTX失败: {e}")
        return None


def read_excel(file_path):
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                text += row_text + '\n'
        return text
    except ImportError:
        print("✗ 未安装 openpyxl")
        return None
    except Exception as e:
        print(f"✗ 读取Excel失败: {e}")
        return None


def read_html(file_path):
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f, 'lxml')
        return soup.get_text(separator='\n')
    except ImportError:
        print("✗ 未安装 beautifulsoup4 或 lxml")
        return None
    except Exception as e:
        print(f"✗ 读取HTML失败: {e}")
        return None


def read_eml(file_path):
    try:
        import email
        from email import policy
        from bs4 import BeautifulSoup
        
        with open(file_path, 'rb') as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
        
        subject = msg.get('Subject', '')
        sender = msg.get('From', '')
        recipient = msg.get('To', '')
        date = msg.get('Date', '')
        
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get("Content-Disposition"))
                
                if "attachment" in content_disposition:
                    continue
                
                if content_type == "text/plain":
                    try:
                        body += part.get_content()
                    except Exception:
                        pass
                elif content_type == "text/html":
                    try:
                        html_content = part.get_content()
                        soup = BeautifulSoup(html_content, 'lxml')
                        body += soup.get_text(separator='\n')
                    except Exception:
                        pass
        else:
            content_type = msg.get_content_type()
            if content_type == "text/plain":
                body = msg.get_content()
            elif content_type == "text/html":
                try:
                    html_content = msg.get_content()
                    soup = BeautifulSoup(html_content, 'lxml')
                    body = soup.get_text(separator='\n')
                except Exception:
                    body = msg.get_content()
        
        content = f"主题: {subject}\n发件人: {sender}\n收件人: {recipient}\n日期: {date}\n\n---\n\n{body}"
        return content
    except Exception as e:
        print(f"✗ 读取EML失败: {e}")
        return None


def read_text_file(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"✗ 读取文件失败: {e}")
        return None


def convert_with_markitdown(input_path, file_type, original_file_path=None):
    temp_md_path = create_temp_file("", suffix='.md')
    
    source_path = original_file_path if original_file_path else input_path
    
    if run_markitdown(input_path, temp_md_path):
        with open(temp_md_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        md_path = create_markdown_with_metadata(content, source_path, file_type)
        delete_temp_file(temp_md_path)
        return md_path
    
    delete_temp_file(temp_md_path)
    return None


def handle_image_file(file_path, file_type):
    print(f"📷 处理图片文件: {file_path}")
    
    ocr_text = extract_text_from_image(file_path)
    
    if not ocr_text:
        print("✗ OCR未能识别到文字")
        return None
    
    temp_txt_path = create_temp_file(ocr_text, suffix='.txt')
    
    md_path = convert_with_markitdown(temp_txt_path, file_type, original_file_path=file_path)
    
    delete_temp_file(temp_txt_path)
    
    if md_path:
        print(f"✅ 图片文件已转换并保存到: {md_path}")
    else:
        print("✗ 图片文件转换失败")
    
    return md_path


def handle_document_file(file_path, file_type):
    print(f"📄 处理文档文件: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()
    
    content = None
    
    if ext == '.pdf':
        content = read_pdf(file_path)
    elif ext in ['.doc', '.docx']:
        content = read_docx(file_path)
    elif ext in ['.ppt', '.pptx']:
        content = read_pptx(file_path)
    elif ext in ['.xlsx', '.xls']:
        content = read_excel(file_path)
    elif ext == '.html':
        content = read_html(file_path)
    elif ext == '.eml':
        content = read_eml(file_path)
    elif ext in TEXT_EXTENSIONS:
        content = read_text_file(file_path)
    
    if content is None:
        print("✗ 无法读取文件内容，尝试使用markitdown直接转换")
        return convert_with_markitdown(file_path, file_type)
    
    temp_txt_path = create_temp_file(content, suffix='.txt')
    
    md_path = convert_with_markitdown(temp_txt_path, file_type, original_file_path=file_path)
    
    delete_temp_file(temp_txt_path)
    
    if md_path:
        print(f"✅ 文档文件已转换并保存到: {md_path}")
    else:
        print("✗ 文档文件转换失败")
    
    return md_path


def run(file_path, file_type):
    if not os.path.exists(file_path):
        print(f"✗ 文件不存在: {file_path}")
        return False
    
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext in IMAGE_EXTENSIONS:
        return handle_image_file(file_path, file_type)
    else:
        return handle_document_file(file_path, file_type)


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("用法: python read.py {文件路径} {文件类型}")
        sys.exit(1)
    
    file_path = sys.argv[1]
    file_type = sys.argv[2]
    run(file_path, file_type)